import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import re
from datetime import datetime
from firebase.config import TEMP_DIR
import os

# Enhanced static theme styling - colors are now hardcoded in the function
THEME_COLORS = {
    "primary": "#8b0000",      # Deep red
    "secondary": "#fcf2f2",    # Light beige
    "accent": "#FFD700",       # Gold
    "text": "#010000",         # Near black
    "background": "#faf7f7",   # Very light gray
    "heading": "#070000",      # Dark black
    "table_header": "#8b0000", # Deep red for table headers
    "table_alt": "#f9f9f9",   # Light gray for alternating rows
    "shadow": "#cccccc"        # Light gray for shadows
}

THEME_FONTS = {
    "heading": "Georgia",
    "body": "Open Sans",
    "code": "Fira Code"
}

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple for pptx."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def apply_theme_font(run, font_type="body", size=None, bold=False, color=None):
    """Apply enhanced theme font settings to a pptx paragraph run."""
    if font_type == "heading":
        run.font.name = THEME_FONTS["heading"]
        run.font.size = size or Pt(32)
        run.font.bold = bold or True
        run.font.color.rgb = hex_to_rgb(color or THEME_COLORS["heading"])
    elif font_type == "code":
        run.font.name = THEME_FONTS["code"]
        run.font.size = size or Pt(14)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color or THEME_COLORS["accent"])
    else:  # body font
        run.font.name = THEME_FONTS["body"]
        run.font.size = size or Pt(18)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color or THEME_COLORS["text"])

def add_enhanced_footer(slide, prs, text, presentation_metadata):
    """Add a styled footer with separator line."""
    # Add separator line
    line_left = Inches(0.5)
    line_top = prs.slide_height - Inches(0.7)
    line_width = prs.slide_width - Inches(1)
    line_height = Inches(0.02)

    line_shape = slide.shapes.add_textbox(line_left, line_top, line_width, line_height)
    line_fill = line_shape.fill
    line_fill.solid()
    line_fill.fore_color.rgb = hex_to_rgb(THEME_COLORS["primary"])

    # Add footer text
    left = Inches(0.5)
    width = Inches(9)
    top = prs.slide_height - Inches(0.5)
    height = Inches(0.4)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()

    footer_text = text.format(presentation=presentation_metadata.get("title", "Presentation"))
    run.text = footer_text
    apply_theme_font(run, font_type="body", size=Pt(10), color=THEME_COLORS["primary"])
    p.alignment = PP_ALIGN.LEFT

def add_slide_number(slide, prs, num):
    """Add styled slide number."""
    left = prs.slide_width - Inches(1.2)
    top = prs.slide_height - Inches(0.5)
    width = Inches(1)
    height = Inches(0.4)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(num)
    apply_theme_font(run, font_type="body", size=Pt(10), bold=True, color=THEME_COLORS["primary"])
    p.alignment = PP_ALIGN.RIGHT

def set_slide_background(slide, bg_hex=None):
    """Set slide background with fallback to theme default."""
    bg = slide.background
    fill = bg.fill
    fill.solid()

    # Use provided color or default theme background
    if bg_hex and re.match(r'^#(?:[0-9a-fA-F]{3}){1,2}$', bg_hex):
        fill.fore_color.rgb = hex_to_rgb(bg_hex)
    else:
        fill.fore_color.rgb = hex_to_rgb(THEME_COLORS["background"])

def add_enhanced_title(slide, title_text, y_offset=Inches(0.4)):
    """Add enhanced title with styling and spacing."""
    title_box = slide.shapes.add_textbox(Inches(0.5), y_offset, Inches(11), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True  # Enable word wrapping
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Auto-resize height to fit text
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    apply_theme_font(run, font_type="heading", size=Pt(24), bold=True)  # Further reduced font size
    p.alignment = PP_ALIGN.LEFT

    return y_offset + Inches(1.2)  # Adjusted spacing

def add_enhanced_subheader(slide, subheader_text, y_offset):
    """Add enhanced subheader with styling."""
    subheader_box = slide.shapes.add_textbox(Inches(0.5), y_offset, Inches(11), Inches(1.0))
    tf = subheader_box.text_frame
    tf.word_wrap = True  # Enable word wrapping
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Auto-resize height to fit text
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subheader_text
    apply_theme_font(run, font_type="body", size=Pt(16), bold=True, color=THEME_COLORS["primary"])
    p.alignment = PP_ALIGN.LEFT

    return y_offset + Inches(0.8)

def add_enhanced_bullets(slide, bullet_points, y_offset):
    """Add enhanced bullet points with better spacing and positioning."""
    # Calculate available space (leave room for footer)
    available_height = Inches(5.5) - y_offset  # Leave more room for footer
    
    bullets_box = slide.shapes.add_textbox(
        Inches(0.8),
        y_offset,
        Inches(10.5),  # Ensure it stays within slide bounds
        max(Inches(1.0), available_height)  # Minimum height with available space
    )

    tf = bullets_box.text_frame
    tf.word_wrap = True  # Enable word wrapping
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Auto-resize to fit content
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    # Clear the default paragraph
    tf.clear()

    for i, bullet in enumerate(bullet_points):
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(4) if i > 0 else Pt(0)
        p.space_after = Pt(4)

        if '**' in bullet:
            parts = bullet.split('**')
            for j, part in enumerate(parts):
                if j % 2 == 0:  # Normal text
                    if part:
                        run = p.add_run()
                        run.text = part
                        apply_theme_font(run, font_type="body", size=Pt(16))
                else:  # Bold text
                    if part:
                        bold_run = p.add_run()
                        bold_run.text = part
                        apply_theme_font(bold_run, font_type="body", size=Pt(16), bold=True)
        else:
            run = p.add_run()
            run.text = bullet
            apply_theme_font(run, font_type="body", size=Pt(16))

    # Return a more conservative y_offset
    return y_offset + min(Inches(3.0), available_height)

def calculate_text_height(text, width, font_size):
    """Calculate approximate height needed for text."""
    # More accurate calculation
    chars_per_line = int(width.inches * 12)  # Approximate characters per line
    lines = max(1, len(text) // chars_per_line + (1 if len(text) % chars_per_line else 0))
    line_height = font_size / 72.0 * 1.2  # Convert points to inches with line spacing
    return Inches(lines * line_height)

def add_enhanced_text_blocks(slide, text_blocks, y_offset):
    """Add enhanced text blocks with better formatting and positioning."""
    current_y = y_offset
    
    for block in text_blocks:
        # Calculate available space
        available_height = Inches(5.5) - current_y  # Leave room for footer
        if available_height <= Inches(0.5):
            break  # Not enough space for more content
            
        width = Inches(10.5)  # Ensure it stays within slide bounds
        
        # Create text box with proper constraints
        tb = slide.shapes.add_textbox(
            Inches(0.8),
            current_y,
            width,
            max(Inches(0.8), available_height - Inches(0.2))  # Minimum height with space buffer
        )

        tf = tb.text_frame
        tf.word_wrap = True  # Enable word wrapping
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Auto-resize to fit content
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)

        p = tf.paragraphs[0]

        if '*' in block and not block.startswith('*Discussion'):
            parts = block.split('*')
            for j, part in enumerate(parts):
                if j % 2 == 0:  # Normal text
                    if part:
                        run = p.add_run()
                        run.text = part
                        apply_theme_font(run, font_type="body", size=Pt(16))
                else:  # Italic text
                    if part:
                        italic_run = p.add_run()
                        italic_run.text = part
                        italic_run.font.italic = True
                        apply_theme_font(italic_run, font_type="body", size=Pt(16))
        else:
            run = p.add_run()
            run.text = block
            apply_theme_font(run, font_type="body", size=Pt(16))

        p.alignment = PP_ALIGN.LEFT
        
        # More conservative spacing calculation
        estimated_lines = len(block) // 80 + 1  # Rough estimate of lines
        block_height = Inches(estimated_lines * 0.25)  # Height per line
        current_y += min(block_height, Inches(1.5)) + Inches(0.15)  # Cap maximum height

    return current_y

def add_enhanced_table(slide, table_json, y_offset):
    """Add enhanced table with better sizing and positioning."""
    if not (table_json.get("headers") and table_json.get("rows")):
        return y_offset

    rows = len(table_json["rows"]) + 1
    cols = len(table_json["headers"])

    # Calculate available space
    available_height = Inches(5.5) - y_offset  # Leave room for footer
    if available_height <= Inches(1):
        return y_offset  # Not enough space for table

    # Calculate table dimensions - ensure it fits within slide bounds
    table_width = Inches(10.5)  # Constrain width to stay within slide
    row_height = Inches(0.35)
    table_height = min(rows * row_height, available_height - Inches(0.3))

    table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), y_offset, table_width, table_height)
    table = table_shape.table

    # Set uniform row height
    for row in table.rows:
        row.height = row_height

    # Style header row
    for c, header in enumerate(table_json["headers"]):
        cell = table.cell(0, c)
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True  # Enable word wrapping in cells
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(header)[:50]  # Limit text length to prevent overflow
        apply_theme_font(run, font_type="heading", size=Pt(14), bold=True, color=THEME_COLORS["secondary"])
        p.alignment = PP_ALIGN.CENTER

        # Header background
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(THEME_COLORS["table_header"])

    # Style data rows
    for r, row in enumerate(table_json["rows"], start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True  # Enable word wrapping in cells
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = tf.add_paragraph()
            
            # Limit cell text length and handle bold formatting
            cell_text = str(val)[:100]  # Limit text length
            
            if '**' in cell_text:
                parts = cell_text.split('**')
                for j, part in enumerate(parts):
                    if j % 2 == 0:  # Normal text
                        if part:
                            run = p.add_run()
                            run.text = part
                            apply_theme_font(run, font_type="body", size=Pt(11))
                    else:  # Bold text
                        if part:
                            bold_run = p.add_run()
                            bold_run.text = part
                            apply_theme_font(bold_run, font_type="body", size=Pt(11), bold=True, color=THEME_COLORS["primary"])
            else:
                run = p.add_run()
                run.text = cell_text
                apply_theme_font(run, font_type="body", size=Pt(11))

            p.alignment = PP_ALIGN.LEFT

            # Alternating row colors
            if r % 2 == 0:
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = hex_to_rgb(THEME_COLORS["table_alt"])

    return y_offset + table_height + Inches(0.3)

def create_enhanced_title_slide(prs, data):
    """Create an enhanced title slide with better styling."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    set_slide_background(slide, THEME_COLORS["background"])

    # Remove default placeholders
    shapes_to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)

    # Enhanced title with proper wrapping
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True  # Enable word wrapping
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Auto-resize to fit content
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data["presentation_metadata"].get("title", "Presentation")
    apply_theme_font(run, font_type="heading", size=Pt(40), bold=True)
    p.alignment = PP_ALIGN.CENTER

    # Enhanced subtitle with proper wrapping
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(1.5))
    tf = subtitle_box.text_frame
    tf.word_wrap = True  # Enable word wrapping
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Auto-resize to fit content
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    subtitle_text = data["presentation_metadata"].get("summary", "")
    if len(subtitle_text) > 150:  # Shorter limit for better wrapping
        subtitle_text = subtitle_text[:150] + "..."
    run.text = subtitle_text
    apply_theme_font(run, font_type="body", size=Pt(16), color=THEME_COLORS["primary"])
    p.alignment = PP_ALIGN.CENTER

    # Author information with proper wrapping
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(0.8))
    tf = author_box.text_frame
    tf.word_wrap = True  # Enable word wrapping
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Presented by {data['presentation_metadata'].get('author', 'AI Education Team')}"
    apply_theme_font(run, font_type="body", size=Pt(14))
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation_from_json(data):
    """Create an enhanced PowerPoint presentation from JSON data."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Create enhanced title slide
    title_slide = create_enhanced_title_slide(prs, data)
    add_enhanced_footer(title_slide, prs, "{presentation} | TeachFlow", data["presentation_metadata"])
    add_slide_number(title_slide, prs, 1)

    # Create content slides
    for idx, slide_json in enumerate(data["slides"], start=2):
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background
        set_slide_background(slide, THEME_COLORS["background"])

        y_offset = Inches(0.3)  # Start higher on the slide

        # Add title
        if slide_json.get("title"):
            y_offset = add_enhanced_title(slide, slide_json["title"], y_offset)

        # Add subheader
        if slide_json["content"].get("subheader"):
            y_offset = add_enhanced_subheader(slide, slide_json["content"]["subheader"], y_offset)

        # Add bullet points
        if slide_json["content"].get("bullet_points"):
            y_offset = add_enhanced_bullets(slide, slide_json["content"]["bullet_points"], y_offset)

        # Add text blocks
        if slide_json["content"].get("text_blocks"):
            y_offset = add_enhanced_text_blocks(slide, slide_json["content"]["text_blocks"], y_offset)

        # Add tables
        for table_json in slide_json["content"].get("tables", []):
            y_offset = add_enhanced_table(slide, table_json, y_offset)

        # Add footer and slide number
        add_enhanced_footer(slide, prs, "{presentation} | TeachFlow", data["presentation_metadata"])
        add_slide_number(slide, prs, idx)

    # Save presentation
    prs_path = os.path.join(TEMP_DIR, f'presentation.pptx')
    prs.save(prs_path)
    print(f"Enhanced presentation saved to {prs_path}")
    return prs_path


def convert_json_to_pptx(result_json):
    cleaned_json_content = str(result_json).strip().removeprefix("```json").removeprefix("```").removesuffix("```").strip()
    return create_presentation_from_json(json.loads(cleaned_json_content))