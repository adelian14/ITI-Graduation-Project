from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from math import ceil

# ==================== 🎨 STYLES ====================
TITLE_FONT_SIZE = Pt(36)
BODY_FONT_SIZE = Pt(20)
TITLE_COLOR = RGBColor(0, 51, 102)         # Navy Blue
BODY_COLOR = RGBColor(30, 30, 30)          # Dark gray
BACKGROUND_COLOR = RGBColor(240, 248, 255) # Light Blue-Gray


# ==================== 🎨 HELPERS ====================
def set_background(slide, color: RGBColor):
    """Set background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def format_textbox(text_frame, font_size=BODY_FONT_SIZE, font_color=BODY_COLOR):
    """Apply font styling to a text frame."""
    for p in text_frame.paragraphs:
        for run in p.runs:
            run.font.size = font_size
            run.font.color.rgb = font_color
        p.alignment = PP_ALIGN.LEFT


# ==================== 📋 AGENDA ====================
def add_agenda_slides(prs, slide_titles, max_per_slide=12):
    """Add agenda slides with two-column layout, split across multiple slides if needed."""
    num_slides = ceil(len(slide_titles) / max_per_slide)

    for i in range(num_slides):
        agenda_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
        set_background(agenda_slide, BACKGROUND_COLOR)
        agenda_slide.shapes.title.text = f"📚 Agenda (Part {i+1})" if num_slides > 1 else "📚 Agenda"

        # Text box area
        left = Inches(0.7)
        top = Inches(1.8)
        width = Inches(8.0)
        height = Inches(5.0)

        textbox = agenda_slide.shapes.add_textbox(left, top, width, height)
        frame = textbox.text_frame
        frame.word_wrap = True
        frame.clear()

        current_chunk = slide_titles[i * max_per_slide : (i + 1) * max_per_slide]
        midpoint = ceil(len(current_chunk) / 2)
        col1 = current_chunk[:midpoint]
        col2 = current_chunk[midpoint:]

        # Column 1
        for idx, topic in enumerate(col1):
            p = frame.add_paragraph()
            p.text = f"{i * max_per_slide + idx + 1}. {topic}"
            p.level = 0

        # Column 2
        for idx, topic in enumerate(col2):
            p = frame.add_paragraph()
            p.text = f"{i * max_per_slide + midpoint + idx + 1}. {topic}"
            p.level = 0
            p.margin_left = Inches(4)

        format_textbox(frame)


# ==================== 🏗️ MAIN FUNCTION ====================
def markdown_to_pptx(slides_dict: dict, output_file: str):
    """
    Generate a beautifully formatted PowerPoint presentation from markdown slide dict.
    Expects each key as: "Session X - Topic Name"
    And value as markdown: Slide 1️⃣ ... **Title:** ... - bullets
    """
    prs = Presentation()

    # Extract clean topic titles for agenda
    topic_titles = [key.split(" - ")[1] if " - " in key else key for key in slides_dict]
    add_agenda_slides(prs, topic_titles)

    for full_key, markdown in slides_dict.items():
        # Divider slide
        session_slide = prs.slides.add_slide(prs.slide_layouts[0])
        set_background(session_slide, BACKGROUND_COLOR)
        session_slide.shapes.title.text = full_key
        try:
            session_slide.placeholders[1].text = "Session Begins"
        except:
            pass

        # Parse individual slides
        for slide_text in markdown.split("Slide ")[1:]:
            slide_lines = slide_text.strip().splitlines()
            slide_title = ""
            bullet_points = []

            for line in slide_lines:
                if line.startswith("**Title:**"):
                    slide_title = line.replace("**Title:**", "").strip()
                elif line.startswith("- "):
                    bullet_points.append(line[2:].strip())

            # Content slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            set_background(slide, RGBColor(255, 255, 255))
            slide.shapes.title.text = slide_title if slide_title else full_key

            # Style title
            title_shape = slide.shapes.title
            p = title_shape.text_frame.paragraphs[0]
            p.runs[0].font.size = TITLE_FONT_SIZE
            p.runs[0].font.color.rgb = TITLE_COLOR

            # Bullet content
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()

            for bullet in bullet_points:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = Pt(4)

            format_textbox(text_frame)

    prs.save(output_file)
    print(f"✅ Presentation saved to {output_file}")
