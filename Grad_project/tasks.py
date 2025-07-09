from crewai import Task
from tqdm import tqdm
from utils.create_pptx import markdown_to_pptx
from tqdm import tqdm
from pptx import Presentation
from pptx.util import Pt
import json 

class ProjectTasks:
    def extract_main_topics(self, agent):
        return Task(
            description=(
                "Generate a list of essential subtopics someone must learn to master a given topic. "
                "Ensure the output is a valid JSON array of strings, without any explanation or extra text."
            ),
            agent=agent,
            expected_output='A valid JSON array of strings representing the key subtopics, e.g., ["Topic A", "Topic B", "Topic C"]'
    )

    def retrieve_combined_context_task(self, agent, vector_store, topics: list, k=30):
    
        topic_context_dict = {}

        for topic in tqdm(topics, desc="🔍 Retrieving Context", unit="topic"):
            try:
                relevant_docs = vector_store.similarity_search(topic, k=k)
                combined_context = "\n\n".join([doc.page_content for doc in relevant_docs])
                topic_context_dict[topic] = combined_context
            except Exception as e:
                print(f"⚠️ Error retrieving context for '{topic}': {e}")
                topic_context_dict[topic] = ""

        return Task(
            description=(
                "For each topic provided, search the vector store to retrieve the most relevant document chunks, "
                "and return the combined text as a dictionary with topic as key and text as value."
            ),
            agent=agent,
            expected_output="A dictionary where each topic maps to a combined string of relevant text chunks.",
            function=lambda: topic_context_dict
        )

    def json_slides_to_pptx(self, slide_dict: dict, output_file: str):
        prs = Presentation()

        for session_title, slides in slide_dict.items():
            # Add a title slide for the session
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = session_title
            try:
                slide.placeholders[1].text = "Session Overview"
            except:
                pass

            for slide_obj in slides:
                title = slide_obj.get("title", "Untitled")
                content = slide_obj.get("content", [])

                ppt_slide = prs.slides.add_slide(prs.slide_layouts[1])
                ppt_slide.shapes.title.text = title

                tf = ppt_slide.placeholders[1].text_frame
                tf.clear()

                for bullet in content:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    p.space_after = Pt(4)

        prs.save(output_file)
        print(f"✅ JSON-based PPTX saved to {output_file}")


    def build_full_course_ppt_task(
        self,
        agent,
        combined_topic_contexts: dict,
        num_weeks: int,
        sessions_per_week: int,
        output_file: str = "Final_Course_Slides.pptx",
        json_output_file: str = "Final_Course_Slides.json"
    ):
        import os

        total_sessions = num_weeks * sessions_per_week
        topics = list(combined_topic_contexts.keys())
        topics_per_session = max(1, len(topics) // total_sessions)

        session_chunks = [
            topics[i:i + topics_per_session]
            for i in range(0, len(topics), topics_per_session)
        ]

        all_slides = {}  # { "Session X - Topic Y": [ {title:..., content:[...]}, ... ] }

        print(f"🗓️ Course Plan: {total_sessions} sessions "
              f"({num_weeks} weeks × {sessions_per_week} sessions/week)")

        for i, session_topics in enumerate(tqdm(session_chunks, desc="🧠 Generating Slides", unit="session")):
            session_number = i + 1
            for topic in session_topics:
                full_context = str(combined_topic_contexts.get(topic, "")).strip()

                if not full_context or len(full_context.split()) < 20:
                    print(f"⚠️ Skipped topic '{topic}' due to empty or too short context.")
                    continue

                prompt = f"""
    You are a professional educational slide creator.

    🎯 Topic: {topic}
    📚 Content: 
    \"\"\"{full_context}\"\"\"

    ✅ Instructions:
    - Create a JSON object with a "slides" array
    - Each slide must be an object with a "title" and a "content" list
    - Each content item is a bullet point (3–5 bullets per slide)
    - Use only the provided context — no outside information
    - Keep tone formal and instructional

    📝 Example:
    {{
      "slides": [
        {{
          "title": "Introduction to ML",
          "content": ["ML is about learning from data.", "It's a core part of AI."]
        }}
      ]
    }}
                """

                try:
                    response = agent.llm.call(prompt)
                    raw_output = response.strip()

                    if not raw_output:
                        print(f"❌ Empty response for topic '{topic}'")
                        continue

                    cleaned_json = (
                        raw_output
                        .removeprefix("```json")
                        .removeprefix("```")
                        .removesuffix("```")
                        .strip()
                    )

                    try:
                        slide_json = json.loads(cleaned_json)
                    except json.JSONDecodeError as e:
                        print(f"❌ JSON decode error for topic '{topic}': {e}\nRaw:\n{cleaned_json[:200]}...")
                        continue

                    if "slides" in slide_json and isinstance(slide_json["slides"], list):
                        slide_key = f"Session {session_number} - {topic}"
                        all_slides[slide_key] = slide_json["slides"]
                    else:
                        print(f"❌ Invalid slide structure for topic '{topic}'")

                except Exception as e:
                    print(f"⚠️ Skipped topic '{topic}' due to error: {e}")

        if all_slides:
            # ✅ Save JSON file
            try:
                with open(json_output_file, 'w', encoding='utf-8') as jf:
                    json.dump(all_slides, jf, indent=2, ensure_ascii=False)
                print(f"✅ JSON saved to: {json_output_file}")
            except Exception as e:
                print(f"❌ Failed to save JSON: {e}")

            # ✅ Save PPTX
            self.json_slides_to_pptx(all_slides, output_file)
            print(f"✅ PowerPoint presentation saved to: {output_file}")
        else:
            print("❌ No slides were generated. Please check your context or prompt.")

        return Task(
            description=(
                "Build a complete multi-session course slide deck from topic-based combined content. "
                "The output includes a PowerPoint presentation and a JSON file of structured slides."
            ),
            agent=agent,
            expected_output=(
                f"A PowerPoint file: {output_file}, and JSON file: {json_output_file}, "
                "containing all course slides structured by session and topic."
            ),
            function=lambda: {"pptx": output_file, "json": json_output_file}
        )
