from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from crewai.tools import tool
import re

@tool
def generate_pptx_from_json(json_data: dict) -> str:
    """
    Generate PowerPoint from structured JSON content, including:
    - Agenda slides with overflow handling
    - Content slides with bullet splitting
    - Code block formatting
    - Spacing, coloring, layout consistency
    """

    # Design & Layout
    primary_color = RGBColor(139, 0, 0)        # Dark red for headers & agenda
    dark_text = RGBColor(30, 30, 30)           # General text
    code_bg = RGBColor(240, 240, 240)          # Code block background
    white_bg = RGBColor(255, 255, 255)         # Background

    font_name = "Calibri"
    code_font_name = "Consolas"
    title_font_size = Pt(36)
    header_font_size = Pt(28)
    content_font_size = Pt(20)
    code_font_size = Pt(16)
    max_bullets_per_slide = 6

    prs = Presentation()

    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    title.text = json_data.get("course_title", "Course Title")
    title.text_frame.paragraphs[0].font.size = title_font_size
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color

    def split_into_bullets(text_or_list):
        """Handles both strings and lists."""
        if isinstance(text_or_list, list):
            return text_or_list
        lines = [p.strip() for p in str(text_or_list).split('\n') if p.strip()]
        bullets = []
        for line in lines:
            if len(line.split()) > 15 or any(x in line for x in [';', '•', '- ', '1.', '2.']):
                split_line = re.split(r'(?<=[.;]) |(?<=\d\.) |(?<=•) |(?<=- )', line)
                bullets.extend([s.strip() for s in split_line if s.strip()])
            else:
                bullets.append(line)
        return bullets

    def add_agenda_slide(agenda_items, part=1):
        """Adds one agenda slide with up to max_bullets_per_slide bullet items."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        suffix = f" (Part {part})" if part > 1 else ""
        title.text = "Agenda" + suffix
        title.text_frame.paragraphs[0].font.size = header_font_size
        title.text_frame.paragraphs[0].font.color.rgb = primary_color
        title.text_frame.paragraphs[0].font.bold = True

        tf = slide.placeholders[1].text_frame
        tf.clear()

        for entry in agenda_items:
            p = tf.add_paragraph()
            p.text = entry['text']
            p.level = entry['level']
            p.font.name = font_name
            p.font.size = header_font_size if entry['level'] == 0 else content_font_size
            p.font.color.rgb = primary_color if entry['level'] == 0 else dark_text
            p.font.bold = entry['level'] == 0
            p.space_after = Pt(4)

    # Generate full agenda list
    full_agenda = []
    for topic in json_data.get("topics", []):
        full_agenda.append({'text': topic.get("topic_title", "Untitled Topic"), 'level': 0})
        for session in topic.get("sessions", []):
            full_agenda.append({'text': session.get("session_title", "Untitled Session"), 'level': 1})

    # Create paginated agenda slides
    i = 0
    part = 1
    while i < len(full_agenda):
        add_agenda_slide(full_agenda[i:i + max_bullets_per_slide], part)
        i += max_bullets_per_slide
        part += 1

    def add_content_slide(session_title, bullets, examples, code, part=1):
        """Adds content slide with header and up to max_bullets_per_slide bullets."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = white_bg

        # Title
        title = slide.shapes.title
        suffix = f" (Part {part})" if part > 1 else ""
        title.text = session_title + suffix
        title.text_frame.paragraphs[0].font.size = header_font_size
        title.text_frame.paragraphs[0].font.color.rgb = primary_color
        title.text_frame.paragraphs[0].font.bold = True

        # Bullet content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        count = 0
        for point in bullets[:max_bullets_per_slide]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.name = font_name
            p.font.size = content_font_size
            p.font.color.rgb = dark_text
            p.space_after = Pt(6)
            count += 1

        # Examples
        for ex in examples:
            for ex_point in split_into_bullets(ex):
                if count >= max_bullets_per_slide:
                    break
                p = tf.add_paragraph()
                p.text = ex_point
                p.level = 0
                p.font.name = font_name
                p.font.size = content_font_size
                p.font.italic = True
                p.font.color.rgb = dark_text
                p.space_after = Pt(6)
                count += 1

        # Code block
        if code and part == 1:
            code_lines = code.strip().split('\n')
            if len(code_lines) > 0:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8), Inches(1.2))
                tf_code = txBox.text_frame
                tf_code.clear()
                p_code = tf_code.add_paragraph()
                p_code.text = code.strip()
                p_code.font.name = code_font_name
                p_code.font.size = code_font_size
                p_code.font.color.rgb = dark_text
                txBox.fill.solid()
                txBox.fill.fore_color.rgb = code_bg

        return bullets[max_bullets_per_slide:]

    # Content Slides
    for topic in json_data.get("topics", []):
        for session in topic.get("sessions", []):
            bullets = split_into_bullets(session.get("content", []))
            examples = session.get("examples", []) if isinstance(session.get("examples", []), list) else []
            code = session.get("code", "")
            part = 1
            while bullets:
                bullets = add_content_slide(session.get("session_title", "Untitled Session"), bullets, examples if part == 1 else [], code if part == 1 else "", part)
                part += 1

    prs.save("course_presentation.pptx")
    return "✅ PowerPoint saved as 'course_presentation.pptx' with multi-page agenda and structured slides."
